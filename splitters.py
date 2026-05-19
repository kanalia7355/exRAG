"""
ファイル拡張子ごとのチャンク分割関数。
新しい形式に対応する場合はこのファイルに分割関数を追加し、
末尾の SPLITTERS 辞書に登録する。

外部ライブラリは関数内で遅延 import しているため、
使わない形式の依存パッケージは入れなくてよい。
"""
from __future__ import annotations

import json as json_lib
from pathlib import Path

from langchain_core.documents import Document
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    MarkdownHeaderTextSplitter,
    RecursiveJsonSplitter,
    Language,
)

import config


# ---------- 共通ユーティリティ ----------

def _read(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _text_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE_TEXT,
        chunk_overlap=config.CHUNK_OVERLAP_TEXT,
        separators=["\n\n", "\n", "。", " ", ""],
    )


def _split_by_language(
    path: Path,
    language: Language,
    filetype: str,
    *,
    text_size: bool = False,
) -> list[Document]:
    """LangChain 組込みの Language 別スプリッタを使う共通処理。"""
    if text_size:
        chunk_size = config.CHUNK_SIZE_TEXT
        overlap = config.CHUNK_OVERLAP_TEXT
    else:
        chunk_size = config.CHUNK_SIZE_CODE
        overlap = config.CHUNK_OVERLAP_CODE

    splitter = RecursiveCharacterTextSplitter.from_language(
        language=language,
        chunk_size=chunk_size,
        chunk_overlap=overlap,
    )
    return [
        Document(
            page_content=c,
            metadata={"source": str(path), "filetype": filetype},
        )
        for c in splitter.split_text(_read(path))
    ]


# ---------- 文章系 ----------

def split_markdown(path: Path) -> list[Document]:
    """Markdown: 見出しで一次分割 → 長い節を日本語向け区切りで二次分割。"""
    text = _read(path)

    header_splitter = MarkdownHeaderTextSplitter(
        headers_to_split_on=[("#", "h1"), ("##", "h2"), ("###", "h3")],
        strip_headers=False,
    )
    docs = header_splitter.split_text(text)

    docs = _text_splitter().split_documents(docs)
    for d in docs:
        d.metadata["source"] = str(path)
        d.metadata["filetype"] = "md"
    return docs


def split_latex(path: Path) -> list[Document]:
    """LaTeX: \\section, \\subsection, 段落などを優先的に区切る。"""
    return _split_by_language(path, Language.LATEX, "tex", text_size=True)


def split_plain_text(path: Path, filetype: str | None = None) -> list[Document]:
    ftype = filetype or path.suffix.lstrip(".")
    return [
        Document(
            page_content=c,
            metadata={"source": str(path), "filetype": ftype},
        )
        for c in _text_splitter().split_text(_read(path))
    ]


# ---------- コード系 ----------

def split_python(path: Path) -> list[Document]:
    return _split_by_language(path, Language.PYTHON, "py")


def split_js(path: Path) -> list[Document]:
    return _split_by_language(path, Language.JS, "js")


def split_html(path: Path) -> list[Document]:
    return _split_by_language(path, Language.HTML, "html", text_size=True)


def split_css(path: Path) -> list[Document]:
    """CSS はルール境界 `}` を最優先の区切り文字に。"""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE_CODE,
        chunk_overlap=config.CHUNK_OVERLAP_CODE,
        separators=["\n}\n", "}\n", "\n\n", "\n", " ", ""],
    )
    return [
        Document(
            page_content=c,
            metadata={"source": str(path), "filetype": "css"},
        )
        for c in splitter.split_text(_read(path))
    ]


# ---------- 構造化データ ----------

def split_json(path: Path) -> list[Document]:
    """JSON: 構造を保ったまま分割。パース不能ならテキスト扱いで fallback。"""
    text = _read(path)
    try:
        data = json_lib.loads(text)
    except Exception:
        return split_plain_text(path, filetype="json")

    splitter = RecursiveJsonSplitter(max_chunk_size=config.CHUNK_SIZE_TEXT)
    chunks = splitter.split_text(json_data=data, convert_lists=True)
    return [
        Document(
            page_content=c,
            metadata={"source": str(path), "filetype": "json"},
        )
        for c in chunks
    ]


# ---------- PDF ----------

def split_pdf(path: Path) -> list[Document]:
    """PDF: ページ単位で読み込み、長いページは2次分割する。
    要 `pypdf` (pip install pypdf)。"""
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(str(path))
    page_docs = loader.load()       # 1 ページ = 1 Document、metadata に page が入る

    docs = _text_splitter().split_documents(page_docs)
    for d in docs:
        d.metadata["source"] = str(path)
        d.metadata["filetype"] = "pdf"
        # page は PyPDFLoader が 0-origin で入れるので 1-origin に直して扱いやすく
        if "page" in d.metadata:
            d.metadata["page"] = int(d.metadata["page"]) + 1
    return docs


# ---------- Word ----------

def split_docx(path: Path) -> list[Document]:
    """DOCX: 全文を取り出して日本語向け区切りで分割。
    要 `docx2txt` (pip install docx2txt)。.doc(旧形式)は非対応。"""
    from langchain_community.document_loaders import Docx2txtLoader

    loader = Docx2txtLoader(str(path))
    raw = loader.load()
    docs = _text_splitter().split_documents(raw)
    for d in docs:
        d.metadata["source"] = str(path)
        d.metadata["filetype"] = "docx"
    return docs


# ---------- PowerPoint ----------

def split_pptx(path: Path) -> list[Document]:
    """PPTX: 1スライド=1 Document として抽出(ノートも含む)。
    長いスライドは2次分割する。
    要 `python-pptx` (pip install python-pptx)。.ppt(旧形式)は非対応。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_docs: list[Document] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        # 各シェイプのテキスト
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    parts.append(line)

        # 発表者ノート
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[ノート]\n{note}")

        if not parts:
            continue

        slide_docs.append(Document(
            page_content="\n".join(parts),
            metadata={
                "source": str(path),
                "filetype": "pptx",
                "slide": idx,
            },
        ))

    return _text_splitter().split_documents(slide_docs)


# ---------- Excel ----------

def split_xlsx(path: Path) -> list[Document]:
    """XLSX: 1シート=1 Document としてタブ区切り形式で抽出。
    長いシートは2次分割する。
    要 `openpyxl` (pip install openpyxl)。.xls(旧形式)は非対応。"""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    sheet_docs: list[Document] = []

    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            if all(v is None for v in row):
                continue
            cells = ["" if v is None else str(v) for v in row]
            rows.append("\t".join(cells))

        if not rows:
            continue

        content = f"# シート: {sheet.title}\n\n" + "\n".join(rows)
        sheet_docs.append(Document(
            page_content=content,
            metadata={
                "source": str(path),
                "filetype": "xlsx",
                "sheet": sheet.title,
            },
        ))

    wb.close()

    # シートが大きい場合に備えて2次分割(タブ区切りなので改行を最優先)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE_TEXT,
        chunk_overlap=config.CHUNK_OVERLAP_TEXT,
        separators=["\n\n", "\n", "\t", " ", ""],
    )
    return splitter.split_documents(sheet_docs)


# ---------- 拡張子 → 分割関数の登録テーブル ----------

SPLITTERS = {
    # 文章
    ".md":   split_markdown,
    ".tex":  split_latex,
    ".txt":  split_plain_text,
    # コード
    ".py":   split_python,
    ".js":   split_js,
    ".html": split_html,
    ".htm":  split_html,
    ".css":  split_css,
    # 構造化
    ".json": split_json,
    # Office / PDF
    ".pdf":  split_pdf,
    ".docx": split_docx,
    ".pptx": split_pptx,
    ".xlsx": split_xlsx,
    # --- 拡張例 ---
    # ".jsx": split_js,
    # ".ts":  lambda p: _split_by_language(p, Language.TS, "ts"),
    # ".tsx": lambda p: _split_by_language(p, Language.TS, "tsx"),
}


def split_file(path: Path) -> list[Document]:
    """登録されていない拡張子はテキスト扱いで分割。"""
    fn = SPLITTERS.get(path.suffix.lower(), split_plain_text)
    return fn(path)
